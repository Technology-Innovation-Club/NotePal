from pypdf import PdfReader
import docx
from pptx import Presentation
import re

# Clean the text
def clean_and_convert_to_utf8(text):
    # Remove any non-ASCII characters and unnecessary whitespaces
    cleaned_text = " ".join(text.split())

    # Convert the cleaned text to UTF-8 encoding
    utf8_text = cleaned_text.encode("utf-8", "ignore").decode("utf-8")

    return utf8_text


# Upload DOCX file
def upload_docx_file(docx_file):
    docx_reader = docx.Document(docx_file)
    return docx_reader


# Get text from DOCX file
def get_docx_text(docx_reader):
    fullText = []
    for para in docx_reader.paragraphs:
        fullText.append(para.text)
    doc_text = "\n".join(fullText)
    cleaned_text = clean_and_convert_to_utf8(doc_text)
    return cleaned_text


# Upload PPTX file
def upload_pptx_file(pptx_file):
    prs = Presentation(pptx_file)
    return prs


# Get text from PPTX file
def get_pptx_text(prs):
    text_runs = []

    # pptx reads in all text for detected shapes in PowerPoint
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    pptx_text = " ".join(text_runs)
    pptx_text_clean = re.sub(r"[^a-zA-Z0-9_.\s]+", "", pptx_text)

    pptx_cleaned_text = clean_and_convert_to_utf8(pptx_text_clean)
    return pptx_cleaned_text


# Upload PDF file
def upload_pdf_file(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    return pdf_reader


# Get text from PDF file
def get_pdf_text(pdf_reader):
    render_page = pdf_reader.pages
    len_reader = len(render_page)
    pdf_text = ""
    for i in range(len_reader):
        page = pdf_reader.pages[i]
        pdf_text += page.extract_text()
    cleaned_text = clean_and_convert_to_utf8(pdf_text)
    return cleaned_text
