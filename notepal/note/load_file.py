from pypdf import PdfReader
import docx
from pptx import Presentation
import pandas as pd
import numpy as np
import re
import os
import datetime


def clean_and_convert_to_utf8(text):
    # Remove any non-ASCII characters and unnecessary whitespaces
    cleaned_text = ' '.join(text.split())
    
    # Convert the cleaned text to UTF-8 encoding
    utf8_text = cleaned_text.encode('utf-8', 'ignore').decode('utf-8')
    
    return utf8_text

# upload DOCX file
def upload_docx_file(docx_file):
    docx_reader = docx.Document(docx_file)
    fullText = []
    for para in docx_reader.paragraphs:
            fullText.append(para.text)
    doc_text = '\n'.join(fullText)
    cleaned_text = clean_and_convert_to_utf8(doc_text)
    return cleaned_text

def pptx_handle_utf8(df_pptx_clean):
    for col in df_pptx_clean.columns:
        if df_pptx_clean[col].dtype==object:
            df_pptx_clean[col] = df_pptx_clean[col].apply(lambda x:np.nan \
                                                          if x==np.nan \
                                                          else str(x).encode('utf-8','replace').decode('utf-8'))

# handle PPTX file
def upload_pptx_file(pptx_file):
    file_dict = {}
    f = open(f'{pptx_file}', "rb")
    prs = Presentation(f)
    text_runs = []

    # pptx reads in all text for detected shapes in powerpoint
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    # append the file name and values to the dictionary
    file_dict[pptx_file] = " ".join(text_runs)
    
    pptx_dict = {}
    # loop through files and pull information
    pptx_dict.update(file_dict)
    
    # creating dataframe from dictionary
    df_pptx = pd.DataFrame([file_dict]).T.reset_index()
    df_pptx.columns = ['file_name','text']

    df_pptx_clean = df_pptx[df_pptx['text'].str.len() > 0].reset_index(drop=True)

    #cleaned text column with all chars except letters, numbers, and periods.
    df_pptx_clean['text_clean'] = df_pptx_clean['text'].apply(lambda x: re.sub(r'[^a-zA-Z0-9_.\s]+','',str(x)))

    # Ensuring UTF-8 encoding
    pptx_handle_utf8(df_pptx_clean)
    
    pptx_text = df_pptx_clean['text_clean'][0]
    
    return pptx_text



# handle PDF file
def upload_pdf_file(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    render_page = pdf_reader.pages
    len_reader = len(render_page)
    pdf_text = ""
    for i in range(len_reader):
        page = pdf_reader.pages[i]
        pdf_text += page.extract_text()
    cleaned_text = clean_and_convert_to_utf8(pdf_text)
    return cleaned_text


def get_file_type(file_path):
    match = re.search(r"[^/\\]+$", file_path)
    if match:
        file_name_type = match.group(0)
        file_type = file_name_type.split('.')[1]
        return file_type
    else:
        return None

def get_file_name(file_path):
    match = re.search(r"[^/\\]+$", file_path)
    if match:
        file_name_type = match.group(0)
        file_name = file_name_type.split('.')[0]
        return file_name
    else:
        return None



def handle_upload(file_path):
    file_type = get_file_type(file_path)
    if file_type == 'pdf':
        text = upload_pdf_file(file_path)
    elif file_type == 'docx':
        text = upload_docx_file(file_path)
    elif file_type == 'pptx':
        text = upload_pptx_file(file_path)
    else:
        text = None
    return text

# Getting the file information
def get_file_current_date(file_path):
    date_time_obj_created = datetime.datetime.fromtimestamp(os.path.getctime(file_path))
    date_created = date_time_obj_created.strftime("%Y-%m-%d")
    return date_created
    
def get_file_modified_date(file_path):
    date_time_obj_modified = datetime.datetime.fromtimestamp(os.path.getmtime(file_path))
    date_modified = date_time_obj_modified.strftime("%Y-%m-%d")
    return date_modified
    
def get_metadata(file_path):
    metadata = {}
    metadata["file_name"] = get_file_name(file_path)
    metadata["file_size"] = os.path.getsize(file_path)
    metadata["file_type"] = get_file_type(file_path)
    metadata["date_created"] = get_file_current_date(file_path)
    metadata["date_modified"] = get_file_modified_date(file_path)
    return metadata

# TEST
# figure out how to use regex to fix this \ problem
file_path = '/code/samples/CSC304-Week-9-Slides.pdf'




print(handle_upload(file_path))
